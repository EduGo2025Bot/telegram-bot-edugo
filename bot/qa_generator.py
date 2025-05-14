# bot/qa_generator.py  –  חילוץ טקסט, GPT עם חיתוך, Cache, מאגר קבוע
import os, re, json, textwrap, hashlib
from pathlib import Path
from typing import List, Dict

# ─────────────  OpenAI (חדש)  ─────────────
from openai import OpenAI

try:
    client = OpenAI()  # Requires env var OPENAI_API_KEY
    _HAS_OPENAI = True
except Exception as e:
    print("⚠️ שגיאה בהתחברות ל־OpenAI:", e)
    _HAS_OPENAI = False

# ─────────────  הגבלות טוקנים/תווים  ─────────────
MAX_CHARS      = 10000
MAX_PAGES_PDF  = 20
MAX_SLIDES_PPT = 20

# ─────────────  מאגר קבוע (bank.json)  ─────────────
def load_bank() -> List[Dict]:
    fp = Path("data/bank.json")
    if fp.exists():
        return json.loads(fp.read_text(encoding="utf-8"))
    return []

def pick_from_bank(k=6):
    import random
    bank = load_bank()
    return random.sample(bank, k=min(k, len(bank))) if bank else _qa_via_placeholder("", k)

# ─────────────  חילוץ טקסט מקובץ  ─────────────
def extract_text(filepath: str) -> str:
    fp = Path(filepath)
    suf = fp.suffix.lower()
    text = ""

    try:
        if suf == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(fp)
            pages = reader.pages[:MAX_PAGES_PDF]
            text = "\n".join(p.extract_text() or "" for p in pages)

        elif suf in {".docx", ".doc"}:
            import docx
            doc = docx.Document(fp)
            text = "\n".join(p.text for p in doc.paragraphs)

        elif suf == ".pptx":
            from pptx import Presentation
            prs = Presentation(fp)
            slides = prs.slides[:MAX_SLIDES_PPT]
            text = "\n".join(
                shape.text
                for slide in slides
                for shape in slide.shapes
                if hasattr(shape, "text") and isinstance(shape.text, str)
            )

    except Exception as e:
        print("❌ שגיאה בקריאת הקובץ:", e)
        return ""

    return text.strip()[:MAX_CHARS]

# ─────────────  Cache לפי hash  ─────────────
CACHE_DIR = "/tmp/qa_cache"
os.makedirs(CACHE_DIR, exist_ok=True)

def _cached(key: str):
    fp = Path(CACHE_DIR) / f"{key}.json"
    if fp.exists():
        t = json.loads(fp.read_text())
        return t
    return None

def _save_cache(key: str, data):
    Path(CACHE_DIR, f"{key}.json").write_text(json.dumps(data, ensure_ascii=False))

# ─────────────  יצירת שאלות  ─────────────
def build_qa_from_text(txt: str, n: int = 6) -> List[Dict]:
    if not isinstance(txt, str):
        print("⚠️ הטקסט שחולץ איננו string אלא:", type(txt))
        return _qa_via_placeholder("", n)

    try:
        key = hashlib.md5(txt.encode()).hexdigest() + f"_{n}"
    except Exception as e:
        print("❌ שגיאה בעיבוד טקסט:", e)
        return _qa_via_placeholder("", n)

    try:
        cached = _cached(key)
        if cached:
            return cached
    except Exception as e:
        print("❌ שגיאה בקריאת cache:", e)
        return _qa_via_placeholder(txt, n)

    try:
        qa = _qa_via_gpt(txt, n) if _HAS_OPENAI else _qa_via_placeholder(txt, n)

        # 🔒 סינון שאלות לא תקינות:
        if isinstance(qa, dict) and "questions" in qa:
            qa["questions"] = [
                q for q in qa["questions"]
                if all(k in q for k in ("question", "options", "type", "correct"))
                and isinstance(q["options"], list)
            ]
        _save_cache(key, qa)
        return qa
    except Exception as e:
        print("❌ שגיאה ביצירת שאלות GPT:", e)
        return _qa_via_placeholder(txt, n)


# --- GPT ---
def _qa_via_gpt(txt: str, n: int):
    prompt = textwrap.dedent(f"""
    צור בדיוק {n} שאלות בעברית על בסיס הטקסט הבא.

    - כל השאלות חייבות להיות במבנה JSON.
    - כל שאלה כוללת את המפתחות הבאים: "question", "type", "options", "correct"
    - סוג השאלה: או "multiple" (5 אפשרויות עם אותיות א.-ה.) או "true_false"
    - החזר בדיוק {n} שאלות בתוך מערך JSON. בלי טקסט נוסף, בלי הסברים.

    הפורמט:
    [
    {{
        "question": "…",
        "type": "multiple",
        "options": ["א. ...", "ב. ...", "ג. ...", "ד. ...", "ה. ..."],
        "correct": "ג"
    }},
    ...
    ]
    """ + txt[:MAX_CHARS])

    rsp = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=2048,
        temperature=0.2,
    )

    content = rsp.choices[0].message.content
    # print("📥 תשובה מ־GPT:\n", content[:300])

    match = re.search(r"\[\s*\{.*?\}\s*\]", content, re.S)
    if not match:
        print("⚠️ לא נמצא JSON תקני בתשובה")
        return _qa_via_placeholder(txt, n)

    try:
        parsed = json.loads(match.group())
        return {"questions": parsed}
    except Exception as e:
        print("❌ שגיאה ב־json.loads():", e)
        return _qa_via_placeholder(txt, n)


# --- Placeholder ---
def _qa_via_placeholder(txt: str, n: int):
    return [
        {
            "type": "true_false",
            "question": "זהו משפט דוגמה – OpenAI לא פעיל.",
            "options": ["נכון", "לא נכון"],
            "correct": "נכון",
        }
    ] * n
